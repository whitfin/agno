# pip install python-pptx

import os
from typing import Any, Dict, List

try:
    from pptx import Presentation
    from pptx.util import Inches
except ImportError:
    raise ImportError("`python-pptx` not installed. Please install using `pip install 'python-pptx'`.")


from agno.tools import Toolkit
from agno.utils.log import logger


class PowerPointTool(Toolkit):
    def __init__(self):
        super().__init__(name="PowerPoint Generator")
        self.register(self.create_ppt)

    def create_ppt(self, slides: List[Dict[str, Any]], filename: str = "presentation.pptx") -> str:
        """Create a PowerPoint presentation from a list of slides.

        Args:
            slides (List[Dict[str, Any]]): A list of slides, each containing:
                - "title" (str): Slide title.
                - "content" (List[str]): Bullet points for the slide.
                - "image_path" (str, optional): Path to an image to include.
            filename (str): The output file name.

        Returns:
            str: Path to the generated PowerPoint file.
        """
        logger.debug(f"printing the slides {slides}")
        prs = Presentation()

        for slide_data in slides:
            slide_layout = prs.slide_layouts[1]  # Title and Content Layout
            slide = prs.slides.add_slide(slide_layout)

            title = slide.shapes.title
            content = slide.placeholders[1]

            title.text = slide_data.get("title", "Untitled Slide")
            for point in slide_data.get("content", []):
                p = content.text_frame.add_paragraph()
                p.text = point

            # Add image if available
            image_path = slide_data.get("image_path")
            if image_path and os.path.exists(image_path):
                left = Inches(1)
                top = Inches(3)
                width = Inches(5)
                slide.shapes.add_picture(image_path, left, top, width=width)

        prs.save(filename)
        return os.path.abspath(filename)
